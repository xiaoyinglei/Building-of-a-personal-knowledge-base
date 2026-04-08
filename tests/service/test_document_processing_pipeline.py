from collections.abc import Callable
from pathlib import Path
from types import SimpleNamespace

import fitz  # type: ignore[import-untyped]
import openpyxl  # type: ignore[import-untyped]
import pytest
from docx import Document as WordDocument
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

import rag.ingest.chunking as document_processing_module
from rag.ingest.chunking import TOCService
from rag.schema.core import ChunkingStrategy, ChunkRole, OcrRegion, OcrResult, SourceType
from tests.support import make_ingest_service


class FakeOcrVisionRepo:
    def extract(self, image_path: Path) -> OcrResult:
        del image_path
        return OcrResult(
            visible_text="Quarterly revenue chart. Figure 1. Revenue grew 20 percent.",
            visual_semantics="bar chart with a caption below it",
            regions=[
                OcrRegion(text="Quarterly revenue chart", bbox=(0, 0, 240, 60)),
                OcrRegion(text="Figure 1", bbox=(0, 61, 120, 90)),
                OcrRegion(text="Revenue grew 20 percent", bbox=(0, 91, 260, 140)),
            ],
        )


class FakeVlmRepo:
    def __init__(self) -> None:
        self.calls = 0

    def describe_visual(self, image_bytes: bytes, *, mime_type: str = "image/png", prompt: str | None = None) -> str:
        self.calls += 1
        assert image_bytes
        assert mime_type == "image/png"
        return "blue chart with a single highlighted rectangle"


def _write_markdown(path: Path) -> None:
    path.write_text(
        "# Quarterly Review\n\n"
        "Overview paragraph for the report.\n\n"
        "## Revenue\n\n"
        "Revenue expanded because enterprise demand improved.\n\n"
        "## Risks\n\n"
        "Supply chain risk remains manageable.\n",
        encoding="utf-8",
    )


def _write_pdf(path: Path) -> None:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Quarterly Review\nRevenue expanded because enterprise demand improved.")
    page = document.new_page()
    page.insert_text((72, 72), "Risks\nSupply chain risk remains manageable.")
    document.save(str(path))
    document.close()


def _write_docx(path: Path, *, with_headings: bool) -> None:
    document = WordDocument()
    if with_headings:
        document.add_heading("Quarterly Review", level=1)
        document.add_paragraph("Overview paragraph for the report.")
        document.add_heading("Revenue", level=2)
        document.add_paragraph("Revenue expanded because enterprise demand improved.")
        document.add_heading("Risks", level=2)
        document.add_paragraph("Supply chain risk remains manageable.")
    else:
        document.add_paragraph("Quarterly Review")
        document.add_paragraph("Overview paragraph for the report.")
        document.add_paragraph("Revenue expanded because enterprise demand improved.")
        document.add_paragraph("Supply chain risk remains manageable.")
    document.save(str(path))


def _write_image(path: Path) -> None:
    image = Image.new("RGB", (640, 320), color="white")
    draw = ImageDraw.Draw(image)
    draw.text((20, 40), "Quarterly revenue chart", fill="black")
    draw.text((20, 120), "Figure 1", fill="black")
    draw.text((20, 200), "Revenue grew 20 percent", fill="black")
    image.save(path)


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Revenue expanded because enterprise demand improved."
    presentation.save(str(path))


def _write_pptx_with_image(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    image_path = path.with_suffix(".png")
    Image.new("RGB", (160, 100), color="blue").save(image_path)
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(2))
    presentation.save(str(path))


def _write_xlsx(path: Path) -> None:
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Revenue"
    sheet["A1"] = "Quarter"
    sheet["B1"] = "Revenue"
    sheet["A2"] = "Q1"
    sheet["B2"] = "120"
    workbook.save(str(path))


@pytest.mark.parametrize(
    ("filename", "writer", "expected_type", "expected_strategy"),
    [
        ("quarterly.md", _write_markdown, SourceType.MARKDOWN, ChunkingStrategy.HIERARCHICAL),
        ("quarterly.pdf", _write_pdf, SourceType.PDF, ChunkingStrategy.HYBRID),
        (
            "quarterly.docx",
            lambda path: _write_docx(path, with_headings=True),
            SourceType.DOCX,
            ChunkingStrategy.HIERARCHICAL,
        ),
        ("quarterly.png", _write_image, SourceType.IMAGE, ChunkingStrategy.IMAGE),
    ],
)
def test_ingest_file_builds_uniform_processing_package(
    tmp_path: Path,
    filename: str,
    writer: Callable[[Path], None],
    expected_type: SourceType,
    expected_strategy: ChunkingStrategy,
) -> None:
    path = tmp_path / filename
    writer(path)
    service = make_ingest_service(tmp_path / "runtime", ocr_repo=FakeOcrVisionRepo())

    result = service.ingest_file(location=str(path), file_path=path, owner="user")

    assert result.processing is not None
    assert result.source.source_type is expected_type
    assert result.processing.analysis.source_type is expected_type
    assert result.processing.routing.selected_strategy is expected_strategy
    assert result.processing.stats.total_chunks == (
        len(result.processing.parent_chunks)
        + len(result.processing.child_chunks)
        + len(result.processing.special_chunks)
    )
    assert result.processing.metadata_summary["schema_version"] == "chunk-pipeline/v1"
    assert result.processing.child_chunks or result.processing.special_chunks


def test_docx_secondary_route_uses_hybrid_when_heading_quality_is_low(tmp_path: Path) -> None:
    path = tmp_path / "low-quality.docx"
    _write_docx(path, with_headings=False)
    service = make_ingest_service(tmp_path / "runtime", ocr_repo=FakeOcrVisionRepo())

    result = service.ingest_file(location=str(path), file_path=path, owner="user")

    assert result.processing is not None
    assert result.processing.analysis.has_dense_structure is False
    assert result.processing.routing.selected_strategy is ChunkingStrategy.HYBRID


def test_image_pipeline_produces_special_chunks_and_parent_child_links(tmp_path: Path) -> None:
    path = tmp_path / "chart.png"
    _write_image(path)
    service = make_ingest_service(tmp_path / "runtime", ocr_repo=FakeOcrVisionRepo())

    result = service.ingest_file(location=str(path), file_path=path, owner="user")

    assert result.processing is not None
    assert {chunk.chunk_role for chunk in result.processing.parent_chunks} == {ChunkRole.PARENT}
    assert all(chunk.parent_chunk_id for chunk in result.processing.child_chunks)
    assert {chunk.special_chunk_type for chunk in result.processing.special_chunks} >= {
        "image_summary",
        "ocr_region",
    }
    assert len({chunk.chunk_id for chunk in result.processing.special_chunks}) == len(result.processing.special_chunks)


def test_ingest_file_supports_pptx_content(tmp_path: Path) -> None:
    path = tmp_path / "quarterly.pptx"
    _write_pptx(path)
    service = make_ingest_service(tmp_path / "runtime", ocr_repo=FakeOcrVisionRepo())

    result = service.ingest_file(location=str(path), file_path=path, owner="user")

    assert result.source.source_type is SourceType.PPTX
    assert result.document.title == "Quarterly Review"
    assert result.visible_text
    assert any("Quarterly Review" in chunk.text for chunk in result.chunks)


def test_ingest_file_supports_xlsx_content(tmp_path: Path) -> None:
    path = tmp_path / "quarterly.xlsx"
    _write_xlsx(path)
    service = make_ingest_service(tmp_path / "runtime", ocr_repo=FakeOcrVisionRepo())

    result = service.ingest_file(location=str(path), file_path=path, owner="user")

    assert result.source.source_type is SourceType.XLSX
    assert result.visible_text
    assert any("Quarter" in chunk.text for chunk in result.chunks)
    assert any("Revenue" in chunk.text for chunk in result.chunks)


def test_ingest_file_enriches_embedded_figure_with_optional_vlm_description(tmp_path: Path) -> None:
    path = tmp_path / "visual-slide.pptx"
    _write_pptx_with_image(path)
    vlm_repo = FakeVlmRepo()
    service = make_ingest_service(
        tmp_path / "runtime",
        ocr_repo=FakeOcrVisionRepo(),
        vlm_repo=vlm_repo,
    )

    result = service.ingest_file(location=str(path), file_path=path, owner="user")

    assert vlm_repo.calls >= 1
    assert "blue chart with a single highlighted rectangle" in result.visible_text
    assert any("blue chart with a single highlighted rectangle" in chunk.text for chunk in result.chunks)


def test_pipeline_reports_clear_error_for_unsupported_file_types(tmp_path: Path) -> None:
    path = tmp_path / "notes.bin"
    path.write_text("not supported", encoding="utf-8")
    service = make_ingest_service(tmp_path / "runtime", ocr_repo=FakeOcrVisionRepo())

    with pytest.raises(ValueError, match="Unsupported file type"):
        service.ingest_file(location=str(path), file_path=path, owner="user")


def test_document_processing_service_uses_cached_tokenizer_configuration(monkeypatch: pytest.MonkeyPatch) -> None:
    captured: dict[str, object] = {}

    def fake_from_pretrained(
        model_name: str,
        *,
        max_tokens: int | None = None,
        local_files_only: bool | None = None,
        **_: object,
    ) -> object:
        captured["model_name"] = model_name
        captured["max_tokens"] = max_tokens
        captured["local_files_only"] = local_files_only
        return object()

    class FakeHybridChunker:
        def __init__(self, *, tokenizer: object, **_: object) -> None:
            captured["tokenizer"] = tokenizer

    monkeypatch.setattr(
        document_processing_module,
        "HuggingFaceTokenizer",
        SimpleNamespace(from_pretrained=fake_from_pretrained),
    )
    monkeypatch.setattr(document_processing_module, "HybridChunker", FakeHybridChunker)

    document_processing_module.DocumentProcessingService(toc_service=TOCService())

    assert captured["model_name"] == "BAAI/bge-m3"
    assert captured["max_tokens"] == 512
    assert captured["local_files_only"] is True
    assert captured["tokenizer"] is not None
